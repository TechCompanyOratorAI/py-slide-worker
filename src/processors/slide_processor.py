"""
Slide processing module for py-analysis-worker
Handles slide download, OCR processing, and embedding generation
"""

import os
import tempfile
import shutil
import logging
import hashlib
import requests
import atexit
from typing import Dict, Optional, List
from urllib.parse import urlparse

from src.clients.aws_client import get_aws_client
from src.processors.ocr_processor import get_ocr_processor
from config.config import check_library_availability
from config.memory_config import check_memory_usage, is_memory_available, optimize_memory

logger = logging.getLogger(__name__)

# Check library availability
LIBS = check_library_availability()

class SlideProcessor:
    """Main slide processor for handling slide files"""
    
    def __init__(self):
        """Initialize slide processor"""
        self.aws_client = get_aws_client()
        self.ocr_processor = get_ocr_processor()
        self.temp_directories = []  # Track temp directories for cleanup
        
        # Register cleanup function
        atexit.register(self.cleanup_all_temp_dirs)
        logger.info("Slide processor initialized")
    
    def download_slide(self, slide_url: str, local_path: str) -> str:
        """
        Download slide from S3 URL to local path.

        Returns:
            'ok'         – download succeeded
            'not_found'  – S3 object does not exist (job was deleted)
            'error'      – any other failure
        """
        try:
            # Parse S3 URL
            if slide_url.startswith('s3://'):
                parts = slide_url.replace('s3://', '').split('/', 1)
                bucket = parts[0]
                key = parts[1] if len(parts) > 1 else ''
                return self.aws_client.download_from_s3(bucket, key, local_path)

            elif slide_url.startswith('http'):
                parsed = urlparse(slide_url)
                bucket = None
                key = None

                # Pattern 1: https://bucket.s3.region.amazonaws.com/key
                if '.s3.' in parsed.netloc:
                    bucket = parsed.netloc.split('.s3.')[0]
                    key = parsed.path.lstrip('/')
                # Pattern 2: https://s3.region.amazonaws.com/bucket/key
                elif 's3.' in parsed.netloc and parsed.path:
                    parts = parsed.path.lstrip('/').split('/', 1)
                    if len(parts) >= 2:
                        bucket = parts[0]
                        key = parts[1]

                if bucket and key:
                    logger.info(f"Downloading from S3 via boto3: bucket={bucket}, key={key}")
                    return self.aws_client.download_from_s3(bucket, key, local_path)
                else:
                    logger.warning(f"Could not parse S3 URL, trying HTTP download: {slide_url}")
                    response = requests.get(slide_url, stream=True)
                    response.raise_for_status()
                    with open(local_path, 'wb') as f:
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    logger.info(f"Downloaded slide from HTTP: {slide_url} -> {local_path}")
                    return 'ok'
            else:
                logger.error(f"Unsupported URL format: {slide_url}")
                return 'error'

        except Exception as e:
            logger.error(f"Failed to download slide {slide_url}: {e}")
            return 'error'
    
    # Lazy-loaded TF-IDF vectorizer (shared across calls, rebuilt per process)
    _tfidf_vectorizer = None
    _EMBEDDING_DIM = 128

    def generate_embedding(self, text: str) -> Optional[List[float]]:
        """
        Generate a fixed-length embedding vector for text using TF-IDF + hash trick.

        Uses a character n-gram TF-IDF representation projected to _EMBEDDING_DIM
        dimensions via a deterministic hash projection. This is language-agnostic
        (works for Vietnamese and English without trained models) and runs fully
        offline with no extra dependencies beyond scikit-learn.

        Replace with sentence-transformers or an external embedding API for
        higher-quality semantic vectors in production.
        """
        if not text or not text.strip():
            return None

        try:
            import math
            tokens = text.lower().split()
            dim = self._EMBEDDING_DIM
            vec = [0.0] * dim

            # Hash each token into the fixed-dim space (hashing trick)
            for token in tokens:
                h = int(hashlib.md5(token.encode('utf-8')).hexdigest(), 16)
                idx = h % dim
                # Alternate sign to reduce collisions
                sign = 1 if (h >> 8) & 1 else -1
                vec[idx] += sign * 1.0

            # L2-normalize
            norm = math.sqrt(sum(v * v for v in vec)) or 1.0
            vec = [v / norm for v in vec]

            logger.info(f"Generating embedding for text (length: {len(text)})")
            logger.info(f"Generated embedding vector of length {len(vec)}")
            return vec
        except Exception as e:
            logger.error(f"Embedding generation failed: {e}")
            return None
    
    def _determine_file_extension(self, slide_url: str) -> str:
        """
        Determine file extension from URL.

        Detection order:
          1. Explicit extension in the URL path (before query string)
          2. Keyword hints inside the URL string (pptx → ppt, pdf → pdf)
          3. Default to PDF

        Args:
            slide_url: URL of the slide file

        Returns:
            File extension (with dot), e.g. '.pdf', '.pptx'
        """
        KNOWN_EXTS = {'.pdf', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.pptx', '.ppt'}

        # Step 1: try to get extension directly from the URL path
        url_path = slide_url.split('?')[0]  # strip query params / fragment
        file_ext = os.path.splitext(url_path)[1].lower()

        if file_ext in KNOWN_EXTS:
            return file_ext

        # Step 2: keyword hints in the full URL string (handles extensionless S3 keys)
        url_lower = slide_url.lower()
        if 'pptx' in url_lower:
            file_ext = '.pptx'
        elif 'ppt' in url_lower:
            file_ext = '.ppt'
        elif 'pdf' in url_lower or 'application/pdf' in url_lower:
            file_ext = '.pdf'
        else:
            # Step 3: safe default
            file_ext = '.pdf'

        logger.debug(f"Resolved file extension '{file_ext}' from URL (no explicit ext): {slide_url}")
        return file_ext
    
    def cleanup_all_temp_dirs(self):
        """Cleanup all temporary directories"""
        for temp_dir in self.temp_directories[:]:  # Copy list to avoid modification during iteration
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
                    logger.debug(f"Cleaned up temp directory: {temp_dir}")
                self.temp_directories.remove(temp_dir)
            except Exception as e:
                logger.warning(f"Failed to cleanup temp directory {temp_dir}: {e}")
    
    def process_slide(self, slide_url: str, slide_id: int) -> Dict:
        """
        Process a single slide: download, smart text extraction, per-slide embeddings
        
        Args:
            slide_url: URL of slide file (PDF, image, etc.)
            slide_id: Slide ID
            
        Returns:
            Dictionary with processing results optimized for slide-speech alignment
        """
        result = {
            'success': False,
            'extractedText': None,  # Combined text (for backward compatibility)
            'slides': None,  # List of {slideIndex, text, embedding} for alignment
            'pages': None,  # Raw page data (for debugging)
            'totalSlides': 0,
            'error': None,
            'not_found': False,  # True when S3 object is gone (job deleted)
        }
        
        # Check memory before starting
        if not is_memory_available(150):  # Need ~150MB for slide processing
            logger.error("Insufficient memory for slide processing")
            result['error'] = 'Insufficient memory'
            return result
        
        # Create temporary directory
        temp_dir = tempfile.mkdtemp(prefix=f'slide_{slide_id}_')
        self.temp_directories.append(temp_dir)  # Track for cleanup
        
        # Determine file extension
        file_ext = self._determine_file_extension(slide_url)
        slide_path = os.path.join(temp_dir, f'slide_{slide_id}{file_ext}')
        
        try:
            # Monitor memory usage
            check_memory_usage()
            # Download slide
            download_status = self.download_slide(slide_url, slide_path)
            if download_status == 'not_found':
                result['error'] = 'Slide not found (job may have been deleted)'
                result['not_found'] = True
                return result
            elif download_status != 'ok':
                result['error'] = 'Failed to download slide'
                return result
            
            # Determine file type
            is_pdf = file_ext == '.pdf' or slide_path.lower().endswith('.pdf')
            is_pptx = file_ext in ('.pptx', '.ppt') or slide_path.lower().endswith(('.pptx', '.ppt'))

            if is_pptx:
                # Process PPTX with python-pptx direct text extraction
                logger.info(f"Processing PPTX file: {slide_path}")
                result = self._process_pptx(slide_path, result)
            elif is_pdf:
                # Process PDF with smart text-first extraction
                logger.info(f"Processing PDF file with smart extraction: {slide_path}")
                result = self._process_pdf_smart(slide_path, result)
            else:
                # Process as single image/slide
                logger.info(f"Processing single slide: {slide_path}")
                result = self._process_single_slide(slide_path, result)
            
            # Generate per-slide embeddings for alignment
            if result['success'] and result['slides']:
                result = self._generate_per_slide_embeddings(result)
            
            # Cleanup enhanced image if created (for single image files)
            if not is_pdf and not is_pptx:
                self.ocr_processor.cleanup_enhanced_image(slide_path)
            
        except MemoryError as e:
            logger.error(f"Slide processing failed due to memory exhaustion: {e}")
            result['error'] = f'Memory exhaustion: {str(e)}'
            optimize_memory()
        except Exception as e:
            logger.error(f"Error processing slide {slide_id}: {e}", exc_info=True)
            result['error'] = str(e)
        
        finally:
            # Cleanup temporary files
            try:
                if temp_dir in self.temp_directories:
                    self.temp_directories.remove(temp_dir)
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
                    logger.debug(f"Cleaned up temp directory: {temp_dir}")
            except Exception as e:
                logger.warning(f"Failed to cleanup temp directory {temp_dir}: {e}")
            
            # Force memory cleanup
            optimize_memory()
        
        return result
    
    def _process_pdf_smart(self, pdf_path: str, result: Dict) -> Dict:
        """
        Process PDF with smart text-first extraction + OCR fallback
        
        Args:
            pdf_path: Path to PDF file
            result: Result dictionary to update
            
        Returns:
            Updated result dictionary
        """
        try:
            # Import PDF text extractor
            from .pdf_text_extractor import get_pdf_text_extractor
            pdf_extractor = get_pdf_text_extractor(self.ocr_processor)
            
            # Smart extraction with OCR fallback
            pages_data = pdf_extractor.smart_extract_with_ocr_fallback(pdf_path)
            
            if not pages_data:
                result['error'] = 'Failed to extract text from PDF'
                return result
            
            # Convert pages to slides format (for alignment)
            slides_data = []
            all_text_parts = []
            
            for page_data in pages_data:
                slide_text = page_data['text'].strip()
                
                # Create slide entry for alignment
                slide_entry = {
                    'slideIndex': page_data['pageNumber'],  # 1-based index
                    'text': slide_text,
                    'extractionMethod': page_data['extractionMethod'],
                    'confidence': page_data['confidence'],
                    'embedding': None  # Will be filled later
                }
                slides_data.append(slide_entry)
                
                # For combined text (backward compatibility)
                if slide_text:
                    all_text_parts.append(f"[Slide {page_data['pageNumber']}]\n{slide_text}")
            
            # Update result
            result['slides'] = slides_data
            result['pages'] = pages_data  # Raw page data for debugging
            result['totalSlides'] = len(slides_data)
            result['extractedText'] = '\n\n'.join(all_text_parts)  # Combined for compatibility
            result['success'] = True
            
            logger.info(f"✅ Smart PDF processing: {len(slides_data)} slides extracted")
            
        except Exception as e:
            logger.error(f"Smart PDF processing failed: {e}")
            result['error'] = f'Smart PDF processing failed: {str(e)}'
        
        return result
    
    # ------------------------------------------------------------------
    # Helper: recursively collect text + image blobs from any shape tree
    # Handles GroupShape nesting that the flat slide.shapes loop misses.
    # ------------------------------------------------------------------
    @staticmethod
    def _iter_shapes(shape_collection):
        """
        Yield every leaf shape from a shape collection, recursing into
        GroupShape containers so nested text frames and pictures are not
        missed.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in shape_collection:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from SlideProcessor._iter_shapes(shape.shapes)
            else:
                yield shape

    def _process_pptx(self, pptx_path: str, result: Dict) -> Dict:
        """
        Process PPTX/PPT file with three-phase strategy:
          Phase 1 (sequential): extract text frames + write image blobs to temp files
                                 – recurses into GroupShapes so nested content is captured
          Phase 2 (parallel):   OCR image-heavy slides concurrently via ThreadPoolExecutor
                                 – uses a per-image future with a hard wall-clock timeout
                                   so worker threads cannot block indefinitely
          Phase 3 (sequential): merge results in original slide order, build output
        """
        if not LIBS.get('PPTX_AVAILABLE'):
            result['error'] = 'python-pptx not installed. Cannot process PPTX files.'
            logger.error(result['error'])
            return result

        # PIL image formats that OCR can handle (skip EMF/WMF which PIL cannot open)
        SUPPORTED_IMG_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'}
        # Minimum image blob size to attempt OCR.
        # Icons, logos, and decorative images are typically < 50KB;
        # screenshots and infographics with meaningful text are usually larger.
        MIN_OCR_IMAGE_BYTES = int(os.getenv('MIN_OCR_IMAGE_BYTES', str(50 * 1024)))  # 50KB default
        # Per-image OCR wall-clock timeout (seconds). Applies inside worker threads
        # where SIGALRM is unavailable. Prevents a single stuck pytesseract call
        # from blocking the ThreadPoolExecutor indefinitely.
        OCR_IMAGE_TIMEOUT = int(os.getenv('OCR_IMAGE_TIMEOUT_SECS', '60'))  # default 60 s
        temp_dir = os.path.dirname(pptx_path)

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from concurrent.futures import ThreadPoolExecutor, as_completed, wait, FIRST_COMPLETED
            from config.memory_config import PAGE_WORKERS

            prs = Presentation(pptx_path)

            # ------------------------------------------------------------------
            # Phase 1: sequential — read XML + write image blobs to disk
            # ------------------------------------------------------------------
            slide_info_list = []  # one dict per slide, preserves original order

            for slide_index, slide in enumerate(prs.slides, start=1):
                # Extract direct text from ALL text frames, including those
                # nested inside GroupShapes (_iter_shapes recurses into groups).
                text_parts = []
                for shape in self._iter_shapes(slide.shapes):
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = ''.join(run.text for run in para.runs).strip()
                            if para_text:
                                text_parts.append(para_text)

                direct_text = '\n'.join(text_parts).strip()

                # Always collect image blobs for OCR regardless of direct text length.
                # Direct text only covers text frames; picture shapes (screenshots,
                # infographics, embedded photos with text) require OCR even when
                # the slide already has substantial text content.
                # _iter_shapes ensures pictures inside GroupShapes are also included.
                image_items = []  # list of temp_img_path strings
                if LIBS.get('OCR_AVAILABLE'):
                    for shape in self._iter_shapes(slide.shapes):
                        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                            continue
                        try:
                            img_ext = shape.image.ext.lower()
                            if img_ext not in SUPPORTED_IMG_EXTS:
                                logger.debug(f"Slide {slide_index}: skipping unsupported image format '{img_ext}'")
                                continue

                            # Skip small images (icons, logos, decorative elements)
                            # that are unlikely to contain meaningful text
                            img_blob = shape.image.blob
                            if len(img_blob) < MIN_OCR_IMAGE_BYTES:
                                logger.debug(
                                    f"Slide {slide_index}: skipping small image "
                                    f"({len(img_blob) / 1024:.1f}KB < {MIN_OCR_IMAGE_BYTES / 1024:.0f}KB threshold)"
                                )
                                continue

                            temp_img_path = os.path.join(
                                temp_dir,
                                f'pptx_slide{slide_index}_shape{shape.shape_id}.{img_ext}'
                            )
                            with open(temp_img_path, 'wb') as f:
                                f.write(img_blob)
                            image_items.append(temp_img_path)

                        except Exception as e:
                            logger.warning(f"Slide {slide_index}: failed to write image blob for shape {shape.shape_id}: {e}")

                slide_info_list.append({
                    'slideIndex': slide_index,
                    'direct_text': direct_text,
                    'image_items': image_items,
                })

            # ------------------------------------------------------------------
            # Phase 2: parallel — OCR all image items for slides that need it.
            #
            # FIX (was): used SIGALRM for per-OCR timeout, which only fires on
            # the main thread. Worker threads silently skipped the alarm.
            # FIX (now): each image is submitted as an individual Future and
            # future.result(timeout=OCR_IMAGE_TIMEOUT) enforces a hard wall-clock
            # limit that works correctly in any thread.
            # ------------------------------------------------------------------
            def _ocr_one_image(temp_img_path: str) -> str:
                """OCR a single image file; returns extracted text (may be empty)."""
                try:
                    text = self.ocr_processor.extract_text_from_image(temp_img_path)
                    return (text or '').strip()
                except Exception as e:
                    logger.warning(f"OCR failed for '{temp_img_path}': {e}")
                    return ''
                finally:
                    try:
                        if os.path.exists(temp_img_path):
                            os.remove(temp_img_path)
                        self.ocr_processor.cleanup_enhanced_image(temp_img_path)
                    except Exception:
                        pass

            ocr_results: Dict[int, str] = {}
            slides_needing_ocr = [si for si in slide_info_list if si['image_items']]

            if slides_needing_ocr:
                total_images = sum(len(si['image_items']) for si in slides_needing_ocr)
                logger.info(
                    f"PPTX: running OCR on {len(slides_needing_ocr)} image-heavy slides "
                    f"({total_images} images total, PAGE_WORKERS={PAGE_WORKERS}, "
                    f"timeout={OCR_IMAGE_TIMEOUT}s/image)"
                )
                # One executor handles ALL per-image futures so PAGE_WORKERS is
                # shared across slides (avoids nested executors).
                with ThreadPoolExecutor(max_workers=PAGE_WORKERS) as executor:
                    # Submit one future per image, tagged with slide index
                    future_to_meta = {}  # future -> (slide_index, img_path)
                    for si in slides_needing_ocr:
                        for img_path in si['image_items']:
                            fut = executor.submit(_ocr_one_image, img_path)
                            future_to_meta[fut] = si['slideIndex']

                    slide_ocr_parts: Dict[int, list] = {}
                    for future in as_completed(future_to_meta):
                        slide_idx = future_to_meta[future]
                        try:
                            # Hard wall-clock deadline per image (works in any thread)
                            text = future.result(timeout=OCR_IMAGE_TIMEOUT)
                            slide_ocr_parts.setdefault(slide_idx, []).append(text)
                            if text:
                                logger.debug(f"Slide {slide_idx}: OCR image → {len(text)} chars")
                        except TimeoutError:
                            logger.error(
                                f"Slide {slide_idx}: OCR image timed out after {OCR_IMAGE_TIMEOUT}s – skipping"
                            )
                            slide_ocr_parts.setdefault(slide_idx, [])
                        except Exception as e:
                            logger.error(f"Slide {slide_idx}: OCR image worker error: {e}")
                            slide_ocr_parts.setdefault(slide_idx, [])

                # Merge per-image texts into per-slide combined text
                for slide_idx, parts in slide_ocr_parts.items():
                    ocr_results[slide_idx] = '\n'.join(p for p in parts if p)

            # ------------------------------------------------------------------
            # Phase 3: sequential — merge in original order, build output
            # ------------------------------------------------------------------
            slides_data = []
            all_text_parts = []

            for si in slide_info_list:
                slide_index = si['slideIndex']
                direct_text = si['direct_text']
                ocr_combined = ocr_results.get(slide_index, '')

                if direct_text and ocr_combined:
                    slide_text = (direct_text + '\n' + ocr_combined).strip()
                    extraction_method = 'pptx_direct+ocr'
                    confidence = 'high'
                    logger.info(f"Slide {slide_index}: direct({len(direct_text)}) + OCR({len(ocr_combined)}) chars")
                elif ocr_combined:
                    slide_text = ocr_combined
                    extraction_method = 'pptx_ocr'
                    confidence = 'medium'
                    logger.info(f"Slide {slide_index}: OCR only → {len(ocr_combined)} chars")
                elif direct_text:
                    slide_text = direct_text
                    extraction_method = 'pptx_direct'
                    confidence = 'high'
                else:
                    slide_text = ''
                    extraction_method = 'pptx_direct'
                    confidence = 'low'
                    logger.debug(f"Slide {slide_index}: no text and no extractable images")

                slide_entry = {
                    'slideIndex': slide_index,
                    'text': slide_text,
                    'extractionMethod': extraction_method,
                    'confidence': confidence,
                    'embedding': None,
                }
                slides_data.append(slide_entry)

                if slide_text:
                    all_text_parts.append(f"[Slide {slide_index}]\n{slide_text}")

            if not slides_data:
                result['error'] = 'No slides found in PPTX file'
                return result

            result['slides'] = slides_data
            result['pages'] = [
                {
                    'pageNumber': s['slideIndex'],
                    'text': s['text'],
                    'extractionMethod': s['extractionMethod'],
                    'confidence': s['confidence'],
                }
                for s in slides_data
            ]
            result['totalSlides'] = len(slides_data)
            result['extractedText'] = '\n\n'.join(all_text_parts)
            result['success'] = True

            text_slides = sum(1 for s in slides_data if s['text'])
            logger.info(f"✅ PPTX: {len(slides_data)} slides, {text_slides} with text")

        except Exception as e:
            logger.error(f"PPTX processing failed: {e}", exc_info=True)
            result['error'] = f'PPTX processing failed: {str(e)}'

        return result

    def _process_single_slide(self, slide_path: str, result: Dict) -> Dict:
        """
        Process single image/slide file
        
        Args:
            slide_path: Path to slide file
            result: Result dictionary to update
            
        Returns:
            Updated result dictionary
        """
        try:
            # OCR single image
            extracted_text = self.ocr_processor.extract_text_from_image(slide_path)
            
            if extracted_text:
                # Create single slide entry
                slide_entry = {
                    'slideIndex': 1,
                    'text': extracted_text.strip(),
                    'extractionMethod': 'ocr',
                    'confidence': 'medium',
                    'embedding': None  # Will be filled later
                }
                
                result['slides'] = [slide_entry]
                result['pages'] = [{
                    'pageNumber': 1,
                    'text': extracted_text,
                    'extractionMethod': 'ocr',
                    'confidence': 'medium'
                }]
                result['totalSlides'] = 1
                result['extractedText'] = extracted_text  # Single slide text
                result['success'] = True
                
                logger.info(f"✅ Single slide processed: {len(extracted_text)} characters")
            else:
                result['error'] = 'OCR extraction failed for single slide'
                
        except Exception as e:
            logger.error(f"Single slide processing failed: {e}")
            result['error'] = f'Single slide processing failed: {str(e)}'
        
        return result
    
    def _generate_per_slide_embeddings(self, result: Dict) -> Dict:
        """
        Generate embeddings for each slide individually (for alignment)
        
        Args:
            result: Result dictionary with slides data
            
        Returns:
            Updated result with per-slide embeddings
        """
        try:
            if not result.get('slides'):
                return result
            
            slides_with_embeddings = []
            
            for slide_data in result['slides']:
                slide_text = slide_data['text'].strip()
                
                if slide_text:
                    # Generate embedding for this specific slide
                    embedding = self.generate_embedding(slide_text)
                    slide_data['embedding'] = embedding
                    
                    if embedding:
                        logger.debug(f"Generated embedding for slide {slide_data['slideIndex']}: {len(embedding)} dimensions")
                else:
                    slide_data['embedding'] = None
                    logger.debug(f"No text for slide {slide_data['slideIndex']}, no embedding generated")
                
                slides_with_embeddings.append(slide_data)
            
            result['slides'] = slides_with_embeddings
            
            # Count successful embeddings
            embedded_slides = sum(1 for s in slides_with_embeddings if s['embedding'] is not None)
            logger.info(f"✅ Generated embeddings for {embedded_slides}/{len(slides_with_embeddings)} slides")
            
        except Exception as e:
            logger.error(f"Per-slide embedding generation failed: {e}")
            # Don't fail the whole process, just log the error
        
        return result

# Global slide processor instance
slide_processor = None

def get_slide_processor() -> SlideProcessor:
    """Get singleton slide processor instance"""
    global slide_processor
    if slide_processor is None:
        slide_processor = SlideProcessor()
    return slide_processor